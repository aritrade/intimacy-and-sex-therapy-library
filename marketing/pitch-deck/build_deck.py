#!/usr/bin/env python3
"""
Generates the investor / PM pitch deck for the Intimacy & Sex Therapy Library
as a real `.pptx` file using python-pptx.

Why python-pptx (and not Google Slides / a Notion doc / Marp):
    - PPTX is the format every investor + PM already opens.
    - python-pptx is dependency-light (pure Python on top of OOXML).
    - The deck is regenerable from this file alone — change the script,
      re-run, get a fresh deck. Numbers and copy live in one place so
      they can't drift.

Output: marketing/pitch-deck/Intimacy-and-Sex-Therapy-Library-Pitch.pptx
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---- Brand palette (matches the live site's design tokens) -----------------

INK_900 = RGBColor(0x14, 0x10, 0x0E)  # near-black serif text
INK_700 = RGBColor(0x3A, 0x33, 0x2E)  # body copy
INK_400 = RGBColor(0x8C, 0x84, 0x7F)  # captions
BG_CREAM = RGBColor(0xF6, 0xEE, 0xE2)  # warm cream background
ACCENT = RGBColor(0xC9, 0x6A, 0x4B)  # coral accent
TEAL = RGBColor(0x3F, 0x7E, 0x7A)  # teal
PLUM = RGBColor(0x6B, 0x46, 0x6D)  # plum

# Standard 16:9 widescreen EMU (PowerPoint's native preset). Inches(13.333)
# rounds to 12,191,695 EMU and PowerPoint flags that as "off-grid", showing a
# "found a problem with content" repair dialog. Pinning to the exact EMU
# value the official template uses keeps PowerPoint happy.
WIDTH = Emu(12_192_000)
HEIGHT = Emu(6_858_000)


def _strip_style(shape) -> None:
    """Remove the auto-generated theme <p:style> element from an auto-shape.

    python-pptx attaches a default theme style block to every add_shape() call.
    PowerPoint occasionally flags those style refs ("PowerPoint found a problem
    with content") when the shape uses an explicit solid fill that contradicts
    the theme. Easiest fix: drop the <p:style> element entirely.
    """
    sp = shape._element
    style = sp.find(qn("p:style"))
    if style is not None:
        sp.remove(style)


def _fill(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    _strip_style(shape)


def _text_box(slide, left, top, width, height, text, *,
              size=18, bold=False, color=INK_900, align=PP_ALIGN.LEFT,
              font="Helvetica Neue"):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    # Allow multi-line strings; first line goes into the empty paragraph,
    # subsequent lines get appended as new paragraphs.
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = font
    return tb


def _bg(slide, color: RGBColor = BG_CREAM) -> None:
    """Paint the slide background via the cSld background API.

    Earlier versions of this script added a full-bleed RECTANGLE auto-shape
    instead. That works, but it left an extra Z-ordered shape that
    PowerPoint flagged in the repair dialog on macOS. The cSld <p:bg>
    element is the OOXML-native way to do this.
    """
    cSld = slide._element.find(qn("p:cSld"))
    existing = cSld.find(qn("p:bg"))
    if existing is not None:
        cSld.remove(existing)
    bg_xml = (
        '<p:bg xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"'
        ' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<p:bgPr><a:solidFill><a:srgbClr val="{rgb}"/></a:solidFill>'
        '<a:effectLst/></p:bgPr></p:bg>'
    ).format(rgb=f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
    cSld.insert(0, etree.fromstring(bg_xml))


def _accent_bar(slide, color: RGBColor = ACCENT) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(0.6),
                                 Inches(0.6), Inches(0.08))
    _fill(bar, color)


def _footer(slide, page_no: int, total: int) -> None:
    _text_box(slide, Inches(0.6), Inches(7.05), Inches(6), Inches(0.3),
              "Intimacy & Sex Therapy Library",
              size=9, color=INK_400)
    _text_box(slide, Inches(11.5), Inches(7.05), Inches(1.3), Inches(0.3),
              f"{page_no:02d} / {total:02d}",
              size=9, color=INK_400, align=PP_ALIGN.RIGHT)


# ---- Slide builders --------------------------------------------------------

def slide_title(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _text_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
              "AN EVIDENCE-GROUNDED, CLINICIAN-CURATED",
              size=14, color=ACCENT, bold=True)
    _text_box(slide, Inches(0.8), Inches(2.1), Inches(11.5), Inches(2.2),
              "Intimacy & Sex Therapy Library",
              size=64, bold=True, color=INK_900, font="Georgia")
    _text_box(slide, Inches(0.8), Inches(4.2), Inches(11.5), Inches(1.5),
              "A non-monetised, India-focused reference library that turns\n"
              "peer-reviewed sex-therapy research into plain language adults\n"
              "and their clinicians can actually use.",
              size=22, color=INK_700, font="Georgia")
    _text_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.4),
              "intimacy-and-sex-therapy-library.vercel.app  ·  Built 2026  ·  Pitch v1",
              size=11, color=INK_400)
    _footer(slide, n, total)


def slide_problem(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "THE PROBLEM",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.4),
              "Adults in India can't get accurate sex-health information\n"
              "from the people who are supposed to provide it.",
              size=34, bold=True, color=INK_900, font="Georgia")

    stats = [
        ("82%", "of Indian medical graduates report receiving < 4 hours\nof sex-education training across all of medical school",
         "Source: Indian J Psychiatry, 2019 cross-sectional study"),
        ("70%+", "of urban Indian couples surveyed by Durex (2017) reported\nsexual concerns they had never raised with any clinician",
         "Source: Durex Global Sex Survey, India sample"),
        ("3", "AASECT-certified sex therapists in all of India (population: 1.4B),\nvs. ~2,800 in the US (population: 330M)",
         "Source: AASECT public directory, 2024"),
        ("\u20B9 0", "is what most existing resources actually cost — but they sit\nbehind paywalls, are written in clinical jargon, or come from\nshame-driven cultural sources",
         "Observation from the 2025 competitive landscape scan"),
    ]
    top = Inches(3.05)
    for i, (figure, text, source) in enumerate(stats):
        x = Inches(0.6 + i * 3.13)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, top, Inches(3.0), Inches(3.4))
        card.adjustments[0] = 0.06
        _fill(card, RGBColor(0xFF, 0xFB, 0xF3))
        _text_box(slide, x + Inches(0.2), top + Inches(0.2),
                  Inches(2.7), Inches(0.9),
                  figure, size=42, bold=True, color=ACCENT, font="Georgia")
        _text_box(slide, x + Inches(0.2), top + Inches(1.25),
                  Inches(2.7), Inches(1.6),
                  text, size=11, color=INK_700)
        _text_box(slide, x + Inches(0.2), top + Inches(2.85),
                  Inches(2.7), Inches(0.5),
                  source, size=8, color=INK_400)
    _footer(slide, n, total)


def slide_why_now(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "WHY NOW",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Three secular shifts make this the moment.",
              size=32, bold=True, color=INK_900, font="Georgia")

    points = [
        ("Cultural permission", TEAL,
         "Indian urban audiences (24-45) now consume sex-positive content\n"
         "in English on YouTube and Instagram at scale — and demand sources\n"
         "they can verify against research, not influencers."),
        ("LLM economics", PLUM,
         "Plain-language summarisation that cost \u20B9 5,000 / paper in 2022\n"
         "now costs <\u20B9 5 / paper with open-source 70B models on Groq.\n"
         "A clinician's worth of writing is suddenly automatable for cents."),
        ("Platform tolerance", ACCENT,
         "Meta and YouTube have moved toward allowing clinician-authored,\n"
         "non-sexualised sex-health content. Reach is still throttled — but\n"
         "well-credentialed, citation-heavy material is no longer auto-banned."),
    ]
    top = Inches(2.9)
    for i, (title, color, body) in enumerate(points):
        y = top + Inches(i * 1.3)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.6), y + Inches(0.1),
                                     Inches(0.3), Inches(0.3))
        _fill(dot, color)
        _text_box(slide, Inches(1.1), y, Inches(2.6), Inches(0.5),
                  title, size=18, bold=True, color=INK_900)
        _text_box(slide, Inches(3.8), y, Inches(9), Inches(1.3),
                  body, size=14, color=INK_700)
    _footer(slide, n, total)


def slide_solution(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "THE SOLUTION",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "An autonomous content engine + a clinician-curated catalog.",
              size=28, bold=True, color=INK_900, font="Georgia")

    pillars = [
        ("Discovery", TEAL,
         "Three nightly agents (PubMed,\nCrossref, Open Library) propose\nnew evidence-grounded resources;\nclinicians approve in one click."),
        ("Writing", PLUM,
         "LLMs draft explainers using a\nbrand + clinical + marketing\nplaybook, then self-critique and\nrewrite until quality thresholds pass."),
        ("Review", ACCENT,
         "Two-stage human gate: a sex\ntherapist signs off clinically; an\neditor signs off editorially. No\nauto-publish, ever."),
        ("Distribution", INK_900,
         "Renders to YouTube / Instagram /\nFacebook Reels with consistent\nvoice + portrait; metrics fed back\ninto the writing loop."),
    ]
    top = Inches(3.0)
    for i, (title, color, body) in enumerate(pillars):
        x = Inches(0.6 + i * 3.15)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, top, Inches(3.0), Inches(3.4))
        card.adjustments[0] = 0.06
        _fill(card, RGBColor(0xFF, 0xFB, 0xF3))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.2), top + Inches(0.25),
                                      Inches(1.4), Inches(0.4))
        chip.adjustments[0] = 0.5
        _fill(chip, color)
        _text_box(slide, x + Inches(0.2), top + Inches(0.28),
                  Inches(1.4), Inches(0.4),
                  title, size=12, bold=True,
                  color=RGBColor(0xFF, 0xFB, 0xF3),
                  align=PP_ALIGN.CENTER)
        _text_box(slide, x + Inches(0.2), top + Inches(0.95),
                  Inches(2.7), Inches(2.3),
                  body, size=12, color=INK_700)
    _footer(slide, n, total)


def slide_how_it_works(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "HOW IT WORKS",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "End-to-end pipeline, fully automated except the two human gates.",
              size=24, bold=True, color=INK_900, font="Georgia")

    steps = [
        ("01", "Discover", "Nightly agent scans PubMed + Crossref + Open Library"),
        ("02", "Draft", "Groq Llama-3-70B writes a clinical, brand-aware script"),
        ("03", "Critique", "LLM grades its own work on 4 axes; rewrites if any fall below 7/10"),
        ("04", "Review", "Sex therapist + editor approve via /admin/queue"),
        ("05", "Render", "Remotion on GitHub Actions: stock photos + Edge TTS + captions"),
        ("06", "Publish", "Human clicks publish; we post to IG / YT / FB with retry + audit"),
        ("07", "Measure", "Weekly poll pulls views, likes, follower deltas; feeds /admin/analytics"),
    ]
    top = Inches(2.9)
    h = Inches(0.55)
    for i, (num, title, body) in enumerate(steps):
        y = top + Inches(i * 0.6)
        _text_box(slide, Inches(0.7), y, Inches(0.7), h,
                  num, size=18, bold=True, color=ACCENT, font="Georgia")
        _text_box(slide, Inches(1.6), y, Inches(2.0), h,
                  title, size=16, bold=True, color=INK_900)
        _text_box(slide, Inches(3.8), y, Inches(9), h,
                  body, size=14, color=INK_700)
    _footer(slide, n, total)


def slide_architecture(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "TECHNICAL ARCHITECTURE",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Mac-free. Workstation-free. Serverless-first, with one escape hatch.",
              size=22, bold=True, color=INK_900, font="Georgia")

    rows = [
        ("App + APIs", "Next.js 14 on Vercel (Mumbai region)", TEAL),
        ("Database", "Neon Postgres + Drizzle ORM + Drizzle migrations", PLUM),
        ("Auth", "Auth.js v5 (magic link + Google OAuth) + role-based gates", ACCENT),
        ("LLM", "Groq Llama-3-70B (primary) → Anthropic Claude Sonnet (fallback)", TEAL),
        ("Voice", "Microsoft Edge TTS (en-US-JennyNeural, no API key required)", PLUM),
        ("Video", "Remotion + FFmpeg on GitHub Actions runners", ACCENT),
        ("Storage", "Vercel Blob (Mumbai edge cache) — no local FS dependencies", TEAL),
        ("Schedules", "3 Vercel crons (daily) + 4 GitHub Actions workflows (hourly/nightly)", PLUM),
        ("Publishing", "Meta Graph API v22 (IG / FB Reels) + YouTube Data API v3 OAuth2", ACCENT),
        ("Observability", "Audit log (PII-scrubbed) + Plausible + Vercel Web Analytics", TEAL),
    ]
    top = Inches(2.9)
    h = Inches(0.4)
    for i, (k, v, color) in enumerate(rows):
        y = top + Inches(i * 0.4)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                     Inches(0.7), y + Inches(0.13),
                                     Inches(0.14), Inches(0.14))
        _fill(dot, color)
        _text_box(slide, Inches(1.0), y, Inches(2.4), h,
                  k, size=12, bold=True, color=INK_900)
        _text_box(slide, Inches(3.6), y, Inches(9), h,
                  v, size=12, color=INK_700)
    _footer(slide, n, total)


def slide_market(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "MARKET",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Big TAM, near-zero credible supply.",
              size=32, bold=True, color=INK_900, font="Georgia")

    tiles = [
        ("\u20B9 11,200 Cr",
         "Indian online wellness market",
         "Statista 2024, projected \u20B9 21,000 Cr by 2028"),
        ("280 M",
         "English-literate urban Indians 18-45",
         "MoSPI 2023 enumeration"),
        ("85 M",
         "Indian couples in the prime sexual-concern\ncohort (married, 25-45)",
         "Census + NFHS-5 cross-tab"),
        ("12 M",
         "AASECT-aligned content seekers globally\nin English-speaking diaspora",
         "Google Trends + YouTube Studio benchmarks"),
    ]
    top = Inches(2.9)
    for i, (figure, label, source) in enumerate(tiles):
        x = Inches(0.6 + i * 3.15)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, top, Inches(3.0), Inches(3.4))
        card.adjustments[0] = 0.06
        _fill(card, RGBColor(0xFF, 0xFB, 0xF3))
        _text_box(slide, x + Inches(0.2), top + Inches(0.3),
                  Inches(2.7), Inches(1.1),
                  figure, size=34, bold=True, color=ACCENT, font="Georgia")
        _text_box(slide, x + Inches(0.2), top + Inches(1.55),
                  Inches(2.7), Inches(1.3),
                  label, size=13, bold=True, color=INK_900)
        _text_box(slide, x + Inches(0.2), top + Inches(2.85),
                  Inches(2.7), Inches(0.5),
                  source, size=8, color=INK_400)
    _footer(slide, n, total)


def slide_competition(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "COMPETITIVE LANDSCAPE",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Each existing option fails on at least one of: clinical, local, or free.",
              size=22, bold=True, color=INK_900, font="Georgia")

    rows = [
        ("Healthify / Practo blogs", "Local, free", "Not clinician-authored;\nSEO-driven, not evidence-driven", INK_400),
        ("Mayo Clinic / Cleveland Clinic", "Clinical, free", "US-centric framing; pricing,\ninsurance, and norms don't transfer", INK_400),
        ("Esther Perel / Emily Nagoski", "Clinical, English", "Behind paywalls + courses;\nnot India-contextualised", INK_400),
        ("Reddit / influencer YouTube", "Local, free", "Variable accuracy; no review;\nshame-spiral risk", INK_400),
        ("Intimacy & Sex Therapy Library", "Clinical + Local + Free",
         "Citation-backed, AASECT-aligned,\nIndia-context, plain-language", ACCENT),
    ]
    top = Inches(3.0)
    for i, (name, wins, gap, color) in enumerate(rows):
        y = top + Inches(i * 0.65)
        is_us = i == len(rows) - 1
        bg_color = RGBColor(0xFF, 0xFB, 0xF3) if not is_us else RGBColor(0xFD, 0xEC, 0xDF)
        row = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.6), y, Inches(12.1), Inches(0.6))
        row.adjustments[0] = 0.2
        _fill(row, bg_color)
        _text_box(slide, Inches(0.85), y + Inches(0.12), Inches(3.6), Inches(0.4),
                  name, size=13, bold=True, color=INK_900)
        _text_box(slide, Inches(4.6), y + Inches(0.12), Inches(2.4), Inches(0.4),
                  wins, size=11, color=color)
        _text_box(slide, Inches(7.2), y + Inches(0.05), Inches(5.4), Inches(0.55),
                  gap, size=11, color=INK_700)
    _footer(slide, n, total)


def slide_traction(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "WHAT'S BUILT (TODAY)",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "A working v1 in production, not slideware.",
              size=32, bold=True, color=INK_900, font="Georgia")

    items = [
        ("Live web app", "intimacy-and-sex-therapy-library.vercel.app"),
        ("Catalog", "100+ clinician-vetted resources, allowlist-restricted to 49 tier-1/2 sources"),
        ("AI surfaces", "Sahay (RAG chat) + Companion (longer-form) + Decide (assessment flow)"),
        ("Autonomous content engine", "Daily script gen + render + multi-platform publish + metrics poll"),
        ("Admin console", "/admin/queue, /admin/analytics, /admin/feedback, /admin/subscribers, /admin/users, /admin/proposals"),
        ("Role + permission model", "viewer / clinician / editor / admin with invite-by-email"),
        ("Quality gates", "Adversarial eval harness + Lighthouse a11y CI on every PR"),
        ("Channels live", "YouTube @IntimacySexTherapyLibrary  ·  Instagram  ·  Facebook page"),
    ]
    top = Inches(2.9)
    for i, (k, v) in enumerate(items):
        y = top + Inches(i * 0.5)
        check = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.7), y + Inches(0.08),
                                       Inches(0.3), Inches(0.3))
        _fill(check, TEAL)
        _text_box(slide, Inches(0.7), y + Inches(0.07),
                  Inches(0.3), Inches(0.3),
                  "\u2713", size=14, bold=True,
                  color=RGBColor(0xFF, 0xFB, 0xF3), align=PP_ALIGN.CENTER)
        _text_box(slide, Inches(1.2), y, Inches(3.6), Inches(0.5),
                  k, size=13, bold=True, color=INK_900)
        _text_box(slide, Inches(4.9), y, Inches(7.8), Inches(0.5),
                  v, size=12, color=INK_700)
    _footer(slide, n, total)


def slide_roadmap(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "ROADMAP",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Three horizons, all building on what's already shipped.",
              size=24, bold=True, color=INK_900, font="Georgia")

    quarters = [
        ("Q3 2026", TEAL, "Stabilise",
         "Newsletter (Buttondown live) → first 1k subs.\n"
         "Multilingual: Hindi + Tamil + Bengali via NLLB-200.\n"
         "Clinician marketplace MVP — book directly through the catalog."),
        ("Q4 2026", PLUM, "Grow",
         "AASECT and ISSWSH content partnerships.\n"
         "Mobile-first Companion app on Android (TWA shell, same backend).\n"
         "Anonymous group-therapy waiting rooms inside the platform."),
        ("Q1 2027", ACCENT, "Monetise",
         "Clinician subscription (\u20B9 999 / mo) for patient-share kits.\n"
         "B2B licensing to insurers and EAP providers.\n"
         "Sponsored research summaries (clearly labelled, never paid placements)."),
    ]
    top = Inches(2.9)
    for i, (q, color, title, body) in enumerate(quarters):
        x = Inches(0.6 + i * 4.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, top, Inches(4.0), Inches(3.7))
        card.adjustments[0] = 0.06
        _fill(card, RGBColor(0xFF, 0xFB, 0xF3))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.25), top + Inches(0.25),
                                      Inches(1.5), Inches(0.4))
        chip.adjustments[0] = 0.5
        _fill(chip, color)
        _text_box(slide, x + Inches(0.25), top + Inches(0.28),
                  Inches(1.5), Inches(0.4),
                  q, size=12, bold=True,
                  color=RGBColor(0xFF, 0xFB, 0xF3),
                  align=PP_ALIGN.CENTER)
        _text_box(slide, x + Inches(0.25), top + Inches(0.95),
                  Inches(3.6), Inches(0.6),
                  title, size=20, bold=True, color=INK_900, font="Georgia")
        _text_box(slide, x + Inches(0.25), top + Inches(1.7),
                  Inches(3.6), Inches(2.0),
                  body, size=11, color=INK_700)
    _footer(slide, n, total)


def slide_economics(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "UNIT ECONOMICS",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Marginal cost of one new explainer: under \u20B9 5.",
              size=32, bold=True, color=INK_900, font="Georgia")

    breakdown = [
        ("LLM draft + critique (Groq Llama-3-70B)", "\u20B9 1.20", "8k input + 2k output tokens"),
        ("Edge TTS narration (free)", "\u20B9 0.00", "Microsoft Edge TTS via public endpoint"),
        ("Stock images (Pexels + Pixabay free API)", "\u20B9 0.00", "Allowlist-restricted to clinically appropriate stock"),
        ("Remotion render (GitHub Actions free tier)", "\u20B9 0.00", "~3 min compute / item, 2,000 free min/month"),
        ("Vercel Blob (Mumbai edge cache)", "\u20B9 1.80", "20 MB per long-form video, 6 GB free tier"),
        ("Meta + YouTube publish APIs", "\u20B9 0.00", "Both free for the volumes we operate at"),
        ("Total marginal cost / explainer", "\u20B9 3.00", "Compare to \u20B9 5,000-15,000 for a freelance clinician"),
    ]
    top = Inches(2.9)
    for i, (k, v, note) in enumerate(breakdown):
        y = top + Inches(i * 0.5)
        is_total = i == len(breakdown) - 1
        weight = True if is_total else False
        color = ACCENT if is_total else INK_900
        _text_box(slide, Inches(0.7), y, Inches(6.5), Inches(0.45),
                  k, size=12, bold=weight, color=color)
        _text_box(slide, Inches(7.2), y, Inches(1.8), Inches(0.45),
                  v, size=14, bold=True, color=color, align=PP_ALIGN.RIGHT)
        _text_box(slide, Inches(9.3), y, Inches(3.4), Inches(0.45),
                  note, size=10, color=INK_400)
    _footer(slide, n, total)


def slide_team(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "TEAM + ADVISORS",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.2),
              "Product-led founder, clinician-gated content.",
              size=28, bold=True, color=INK_900, font="Georgia")

    blocks = [
        ("Founder / Engineering", TEAL,
         "Aritra De\nProduct + full-stack engineer.\nBuilt the engine; runs operations.\nFinal call on architecture and brand."),
        ("Clinical review", PLUM,
         "Engages credentialed sex therapists\n(AASECT-aligned standards) for every\npublished explainer. The library does\nnot publish unreviewed material."),
        ("Editorial review", ACCENT,
         "Second-stage human gate before any\npost goes live: factual accuracy,\ntone, India-context, links sanity-checked.\nNever auto-published."),
    ]
    top = Inches(3.0)
    for i, (role, color, body) in enumerate(blocks):
        x = Inches(0.6 + i * 4.2)
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x, top, Inches(4.0), Inches(3.6))
        card.adjustments[0] = 0.06
        _fill(card, RGBColor(0xFF, 0xFB, 0xF3))
        chip = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                      x + Inches(0.25), top + Inches(0.25),
                                      Inches(2.4), Inches(0.4))
        chip.adjustments[0] = 0.5
        _fill(chip, color)
        _text_box(slide, x + Inches(0.25), top + Inches(0.28),
                  Inches(2.4), Inches(0.4),
                  role, size=11, bold=True,
                  color=RGBColor(0xFF, 0xFB, 0xF3),
                  align=PP_ALIGN.CENTER)
        _text_box(slide, x + Inches(0.25), top + Inches(0.95),
                  Inches(3.6), Inches(2.5),
                  body, size=13, color=INK_700)
    _footer(slide, n, total)


def slide_ask(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _accent_bar(slide)
    _text_box(slide, Inches(0.6), Inches(0.85), Inches(12), Inches(0.5),
              "THE ASK",
              size=12, bold=True, color=ACCENT)
    _text_box(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(2.2),
              "We're not asking for capital today.\nWe're asking for two things money can't buy:",
              size=30, bold=True, color=INK_900, font="Georgia")

    asks = [
        ("01 — Clinical advisors",
         "Sex therapists, OB-GYNs, andrologists, couples therapists.\n"
         "Sign two explainers a month, get co-authorship credit and a\n"
         "permanent link to your practice from every page you sign."),
        ("02 — Distribution partnerships",
         "AASECT, ISSWSH, Indian medical colleges, and EAP providers\n"
         "who want a credible asset to give patients between sessions.\n"
         "Co-branding is on the table; revenue isn't required."),
    ]
    top = Inches(4.0)
    for i, (title, body) in enumerate(asks):
        y = top + Inches(i * 1.4)
        _text_box(slide, Inches(0.6), y, Inches(12), Inches(0.5),
                  title, size=18, bold=True, color=ACCENT, font="Georgia")
        _text_box(slide, Inches(0.6), y + Inches(0.5), Inches(12), Inches(1.0),
                  body, size=13, color=INK_700)
    _footer(slide, n, total)


def slide_closing(prs: Presentation, n: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide)
    _text_box(slide, Inches(0.6), Inches(2.6), Inches(12), Inches(1.5),
              "Adults in India deserve\nbetter sex-health information.",
              size=42, bold=True, color=INK_900, font="Georgia",
              align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(0.6), Inches(4.4), Inches(12), Inches(0.6),
              "We're already shipping it. Let's make it reach them.",
              size=20, color=INK_700, font="Georgia",
              align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(0.6), Inches(5.7), Inches(12), Inches(0.6),
              "intimacy-and-sex-therapy-library.vercel.app",
              size=14, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, Inches(0.6), Inches(6.1), Inches(12), Inches(0.4),
              "Pitch contact: aritrajob79@gmail.com",
              size=12, color=INK_400, align=PP_ALIGN.CENTER)
    _footer(slide, n, total)


# ---- Driver ----------------------------------------------------------------

def main() -> None:
    prs = Presentation()
    prs.slide_width = WIDTH
    prs.slide_height = HEIGHT
    # python-pptx's default template is 4:3; setting slide_width/height only
    # changes the dimensions, not the `type` attribute on <p:sldSz>. Without
    # this fixup, PowerPoint sees 16:9 EMU dimensions with type="screen4x3"
    # and shows a "found a problem with content" repair prompt. Switching the
    # type to "screen16x9" matches the dimensions and clears the warning.
    sld_sz = prs.part._element.find(qn("p:sldSz"))
    if sld_sz is not None:
        sld_sz.set("type", "screen16x9")

    builders = [
        slide_title,
        slide_problem,
        slide_why_now,
        slide_solution,
        slide_how_it_works,
        slide_architecture,
        slide_market,
        slide_competition,
        slide_traction,
        slide_economics,
        slide_roadmap,
        slide_team,
        slide_ask,
        slide_closing,
    ]
    total = len(builders)
    for i, build in enumerate(builders, start=1):
        build(prs, i, total)

    out = Path(__file__).resolve().parent / "Intimacy-and-Sex-Therapy-Library-Pitch.pptx"
    prs.save(out)
    print(f"Wrote {out}  ({total} slides)")


if __name__ == "__main__":
    main()
